import streamlit as st
import json
import os
import tempfile
import io
from urllib.parse import quote
from datetime import datetime, timezone, timedelta
from google import genai
from google.genai import types
from prompts import get_sales_prompt
from utils import clean_json_output, convert_currency, safe_parse_json
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors

# Import the PowerPoint generator library
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    pass # Handled via pip install instruction

# ==========================================
# LOCALIZED SALES PDF ENGINE
# ==========================================
def _safe_text(text):
    return str(text).replace("₹", "INR ").encode('ascii', 'ignore').decode('ascii')

def _draw_wrapped_text(c, text, x, y, max_width, font_name="Helvetica", font_size=11, line_height=14):
    c.setFont(font_name, font_size)
    words = _safe_text(text).split(' ')
    line = ""
    for word in words:
        if c.stringWidth(line + word + " ", font_name, font_size) < max_width:
            line += word + " "
        else:
            c.drawString(x, y, line.strip())
            y -= line_height
            line = word + " "
    if line:
        c.drawString(x, y, line.strip())
        y -= line_height
    return y

def _check_page_break(c, current_y, height, threshold=100):
    if current_y < threshold:
        c.showPage()
        return height - 50
    return current_y

def _draw_header(c, width, height, title):
    c.setFillColorRGB(0.133, 0.545, 0.133) 
    c.rect(0, height - 100, width, 100, fill=1, stroke=0)
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 22)
    c.drawString(40, height - 60, title)
    
    ist = timezone(timedelta(hours=5, minutes=30))
    date_str = datetime.now(ist).strftime("%Y-%m-%d %H:%M IST")
    c.setFont("Helvetica", 10)
    c.drawString(40, height - 80, _safe_text(f"Generated on: {date_str} | BridgeBuild AI"))
    if os.path.exists("Logo_bg_removed.png"):
        c.drawImage("Logo_bg_removed.png", width - 100, height - 90, width=80, height=80, mask='auto')

def generate_local_sales_pdf(ticket_data, currency="USD ($)"):
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    
    _draw_header(c, width, height, "Sales & Feasibility Report")
    y = height - 140
    c.setFillColor(colors.black)
    
    y = _draw_wrapped_text(c, "Project Summary:", 40, y, 500, "Helvetica-Bold", 14)
    y -= 5
    y = _draw_wrapped_text(c, ticket_data.get('project_summary', 'Untitled Project'), 40, y, 500, "Helvetica", 11)
    y -= 15

    y = _check_page_break(c, y, height)
    score = ticket_data.get('feasibility_score', 'N/A')
    y = _draw_wrapped_text(c, f"Feasibility Score: {score}", 40, y, 500, "Helvetica-Bold", 12)
    y -= 5
    y = _draw_wrapped_text(c, ticket_data.get('feasibility_reason', ''), 40, y, 500, "Helvetica", 10)
    y -= 15

    raw_budget = ticket_data.get('budget_estimate_usd', '0')
    low = convert_currency(raw_budget.split("-")[0].strip() if "-" in raw_budget else raw_budget, currency)
    high = convert_currency(raw_budget.split("-")[1].strip() if "-" in raw_budget else raw_budget, currency)
    budget = f"{low} - {high}"

    y = _draw_wrapped_text(c, f"Estimated Timeline: {ticket_data.get('estimated_timeline', 'N/A')}", 40, y, 500, "Helvetica-Bold", 12)
    y = _draw_wrapped_text(c, f"Estimated Budget: {budget}", 40, y, 500, "Helvetica-Bold", 12)
    y -= 20

    y = _check_page_break(c, y, height)
    c.setFillColorRGB(0.133, 0.545, 0.133)
    c.drawString(40, y, _safe_text("Critical 'Ask' List for Client:"))
    y -= 20
    c.setFillColor(colors.black)
    
    for q in ticket_data.get("client_questions", []):
        if q.strip(): 
            y = _check_page_break(c, y, height)
            y = _draw_wrapped_text(c, f"• {q}", 40, y, 500, "Helvetica", 11)
            y -= 5

    y -= 10
    y = _check_page_break(c, y, height)
    c.setFillColorRGB(0.7, 0.1, 0.1)
    c.drawString(40, y, _safe_text("Deal Breakers & Risks:"))
    y -= 20
    c.setFillColor(colors.black)
    
    for r in ticket_data.get("deal_breakers", []):
        if r.strip():
            y = _check_page_break(c, y, height)
            y = _draw_wrapped_text(c, f"• {r}", 40, y, 500, "Helvetica", 11)
            y -= 5

    c.save()
    buffer.seek(0)
    return buffer

# ==========================================
# SALES PPTX PITCH DECK ENGINE
# ==========================================
def generate_sales_pptx(ticket_data, currency="USD ($)"):
    """Generates a 5-slide pitch deck dynamically from the Sales JSON data."""
    prs = Presentation()
    
    # 1. TITLE SLIDE
    slide_layout = prs.slide_layouts[0] # 0 is standard Title Slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Project Proposal & Feasibility"
    summary_text = ticket_data.get('project_summary', 'New Project')
    subtitle.text = summary_text[:80] + "..." if len(summary_text) > 80 else summary_text

    # 2. EXECUTIVE SUMMARY SLIDE
    slide_layout = prs.slide_layouts[1] # 1 is Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Executive Summary"
    body.text = ticket_data.get('project_summary', 'No summary provided.')

    # 3. FEASIBILITY & TIMELINE SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Feasibility & Timeline Details"
    tf = body.text_frame
    tf.text = f"Feasibility Assessment: {ticket_data.get('feasibility_score', 'N/A')}"
    
    p = tf.add_paragraph()
    p.text = f"Reasoning: {ticket_data.get('feasibility_reason', 'N/A')}"
    p.space_after = Pt(14)
    
    p2 = tf.add_paragraph()
    p2.text = f"Estimated Timeline: {ticket_data.get('estimated_timeline', 'N/A')}"

    # 4. INVESTMENT (BUDGET) SLIDE
    raw_cost = ticket_data.get("budget_estimate_usd", "0-0")
    low_end = raw_cost.split("-")[0] if "-" in raw_cost else raw_cost
    high_end = raw_cost.split("-")[1] if "-" in raw_cost else raw_cost
    fmt_low = convert_currency(low_end, currency)
    fmt_high = convert_currency(high_end, currency)

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Estimated Investment (MVP)"
    tf = body.text_frame
    tf.text = f"Estimated Budget Scope: {fmt_low} - {fmt_high}"
    p = tf.add_paragraph()
    p.text = "Note: Final scoping requires technical architecture breakdown."

    # 5. NEXT STEPS & CLARIFICATIONS SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Next Steps & Clarifications"
    tf = body.text_frame
    tf.text = "To proceed to the Technical Scoping phase, we need to clarify:"
    
    for q in ticket_data.get("client_questions", []):
        p = tf.add_paragraph()
        p.text = q
        p.level = 1

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


# ==========================================
# SALES DASHBOARD RENDERER
# ==========================================
def render_sales_dashboard(supabase):
    st.title("BridgeBuild AI - Sales Intake")
    st.markdown("### Quickly validate requirements and get estimated timelines.")

    uploaded_file = st.file_uploader("Upload Client Audio (.mp3, .wav)", type=["mp3", "wav", "m4a"])
    
    if uploaded_file:
        file_ext = uploaded_file.name.split('.')[-1].lower()
        if file_ext in ['mp3', 'wav', 'm4a']:
            st.audio(uploaded_file)
            
    sales_input = st.text_area("Or Paste Notes:", height=150, placeholder="Example: Client needs a basic e-commerce site with Stripe integration.")

    api_key = st.secrets.get("GOOGLE_API_KEY")
    
    user_prefs = st.session_state.get("user_prefs", {})
    currency = user_prefs.get("currency", "USD ($)")
    rate_type = user_prefs.get("rate_standard", "US Agency ($150/hr)")
    model_choice = user_prefs.get("ai_model", "Gemini 1.5 Flash (Fast)")

    if "active_sales_ticket" not in st.session_state: st.session_state.active_sales_ticket = None
    if "active_sales_ticket_id" not in st.session_state: st.session_state.active_sales_ticket_id = None

    if st.button("Analyze Request", type="primary"):
        if not api_key:
            st.error("System Error: AI Engine is currently offline.")
        elif not sales_input and not uploaded_file:
            st.warning("Please enter text or upload a file to proceed.")
        else:
            try:
                client = genai.Client(api_key=api_key)
                SALES_PROMPT = get_sales_prompt(rate_type)

                with st.status("Initializing AI Engine...", expanded=True) as status:
                    st.write("Parsing input and extracting requirements...")
                    model_id = "gemini-2.5-flash" if "Flash" in model_choice else "gemini-2.5-pro"
                    prompt_contents = []
                    
                    if uploaded_file:
                        file_ext = f".{uploaded_file.name.split('.')[-1]}"
                        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
                            tmp.write(uploaded_file.getvalue())
                            tmp_path = tmp.name
                        gemini_file = client.files.upload(file=tmp_path)
                        prompt_contents.append(gemini_file)
                        os.remove(tmp_path)
                    
                    st.write("Consulting Sales & Engineering models...")
                    text_instruction = sales_input if sales_input else "Analyze this client request."
                    prompt_contents.append(text_instruction)
                    
                    response = client.models.generate_content(
                        model=model_id, 
                        config=types.GenerateContentConfig(system_instruction=SALES_PROMPT, temperature=0.0, response_mime_type="application/json"),
                        contents=prompt_contents
                    )
                    
                    if uploaded_file:
                        client.files.delete(name=gemini_file.name)
                    
                    st.write("Formatting budget and feasibility metrics...")
                    data, error_msg = safe_parse_json(response.text)
                    
                    if error_msg:
                        status.update(label="Analysis Failed", state="error", expanded=True)
                        st.error(error_msg)
                        st.stop()
                    else:
                        status.update(label="Analysis Complete!", state="complete", expanded=False)
                        st.session_state.active_sales_ticket = data

                raw_cost = data.get("budget_estimate_usd", "0-0")
                low_end = raw_cost.split("-")[0] if "-" in raw_cost else raw_cost
                high_end = raw_cost.split("-")[1] if "-" in raw_cost else raw_cost
                fmt_low = convert_currency(low_end, currency)
                fmt_high = convert_currency(high_end, currency)
                score = data.get("feasibility_score", "Yellow")

                new_ticket = {
                    "user_id": st.session_state.user.id,
                    "summary": data.get("project_summary", "Sales Intake"),
                    "cost": f"{fmt_low} - {fmt_high}",
                    "raw_cost": raw_cost,
                    "complexity": score, 
                    "time": data.get("estimated_timeline", "Unknown"),
                    "full_data": json.dumps(data),
                    "status": "Draft",               
                    "target_department": "None"      
                }
                db_res = supabase.table("tickets").insert(new_ticket).execute()
                st.session_state.active_sales_ticket_id = db_res.data[0]['id']

            except Exception as e:
                st.error(f"An error occurred: {str(e)}")

    # ==========================================
    # RENDER THE ACTIVE SALES UI (EXECUTIVE POLISH)
    # ==========================================
    if st.session_state.active_sales_ticket:
        data = st.session_state.active_sales_ticket
        
        st.write("")
        st.success("Analysis Ready! Review the Executive Summary below.")
        
        # --- EXECUTIVE SUMMARY CARD ---
        with st.container(border=True):
            score = data.get("feasibility_score", "Yellow")
            if "Green" in score: 
                st.markdown("### 🟢 **Feasibility: GREEN (Highly Feasible)**")
                st.success(f"**Reason:** {data.get('feasibility_reason')}")
            elif "Red" in score: 
                st.markdown("### 🔴 **Feasibility: RED (High Risk)**")
                st.error(f"**Warning:** {data.get('feasibility_reason')}")
            else: 
                st.markdown("### 🟡 **Feasibility: YELLOW (Needs Scoping)**")
                st.warning(f"**Reason:** {data.get('feasibility_reason')}")
                
            st.markdown(f"**Project Summary:** {data.get('project_summary')}")

        st.write("")

        # --- BUDGET & TIMELINE METRICS ---
        raw_cost = data.get("budget_estimate_usd", "0-0")
        low_end = raw_cost.split("-")[0] if "-" in raw_cost else raw_cost
        high_end = raw_cost.split("-")[1] if "-" in raw_cost else raw_cost
        fmt_low = convert_currency(low_end, currency)
        fmt_high = convert_currency(high_end, currency)
        
        with st.container(border=True):
            col1, col2 = st.columns(2)
            with col1: st.metric("Estimated MVP Timeline", data.get("estimated_timeline"))
            with col2: st.metric("Estimated MVP Budget", f"{fmt_low} - {fmt_high}")
            
        st.write("")
        
        # --- THE 'ASK' LIST & DEAL BREAKERS ---
        col_q, col_r = st.columns(2)
        with col_q:
            with st.container(border=True):
                st.markdown("#### The 'Ask' List")
                st.caption("Critical missing info to ask the client:")
                for q in data.get("client_questions", []):
                    prefix = q.split(":", 1)[0] + ":" if ":" in q else ""
                    rest = q.split(":", 1)[1] if ":" in q else q
                    st.info(f"**{prefix}** {rest.strip()}" if prefix else rest)
                    
        with col_r:
            with st.container(border=True):
                st.markdown("#### Deal Breakers")
                st.caption("Major risks to evaluate before signing:")
                for r in data.get("deal_breakers", []):
                    prefix = r.split(":", 1)[0] + ":" if ":" in r else ""
                    rest = r.split(":", 1)[1] if ":" in r else r
                    st.error(f"**{prefix}** {rest.strip()}" if prefix else rest)

        st.divider()
        col_action1, col_action2 = st.columns([1, 1], gap="medium")
        
        with col_action1:
            st.markdown("#### Export Sales Artifacts")
            st.download_button("Download Sales PDF", data=generate_local_sales_pdf(data, currency), file_name="bridgebuild_sales_report.pdf", mime="application/pdf", use_container_width=True)
            
            # --- NEW PPTX EXPORT BUTTON ---
            st.download_button("Download Pitch Deck (.pptx)", data=generate_sales_pptx(data, currency), file_name="bridgebuild_pitch_deck.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
            
        with col_action2:
            st.markdown("#### Email Sales Team")
            ticket_name = data.get('project_summary', 'New Project Request')[:50] + "..."
            sales_body = f"Hello Sales Team,\n\nPlease review the initial Sales & Feasibility scoping for the upcoming client request.\n\n-> FEASIBILITY SCORE: {score}\n-> EST. TIMELINE: {data.get('estimated_timeline', 'N/A')}\n-> EST. BUDGET: {fmt_low} - {fmt_high}\n\n-> PROJECT SUMMARY:\n{data.get('project_summary', 'N/A')}\n\n-> CRITICAL 'ASK' LIST (Questions for the Client):\n"
            for q in data.get("client_questions", []): sales_body += f"- {q}\n"
            sales_body += f"\n-> DEAL BREAKERS / RISKS:\n"
            for r in data.get("deal_breakers", []): sales_body += f"- {r}\n"
            sales_body += "\nBest,\nBridgeBuild Sales Hub"
            sales_mailto = f"mailto:?subject={quote(f'Sales Quote: {ticket_name}')}&body={quote(sales_body)}"
            
            st.markdown(f"""
                <a href="{sales_mailto}" target="_blank" style="text-decoration: none;">
                    <button style="width: 100%; background-color: #2E7D32; color: white; border: none; padding: 0.5rem 1rem; border-radius: 0.5rem; font-weight: 400; cursor: pointer; line-height: 1.6; font-size: 1rem; margin-top: 2px;">
                        Email Sales Summary
                    </button>
                </a>
                """, unsafe_allow_html=True)
            
        # ==========================================
        # THE HANDOFF BUTTON!
        # ==========================================
        st.divider()
        st.markdown("#### Department Handoff")
        st.info("Client approved the quote? Send this locked scope directly to the Product Management team to build the Agile Ticket.")
        
        if st.session_state.active_sales_ticket_id:
            if st.button("Approve Quote & Send to PM Hub", type="primary", use_container_width=True):
                try:
                    supabase.table("tickets").update({"status": "Awaiting PM Scoping", "target_department": "PM"}).eq("id", st.session_state.active_sales_ticket_id).execute()
                    st.success("Successfully routed to the PM Inbox.")
                except Exception as e:
                    st.error(f"Failed to handoff ticket: {str(e)}")
            
    # --- SALES HISTORY SECTION (POLISHED) ---
    st.divider()
    st.subheader("Saved Sales Quotes")
    
    try:
        db_response = supabase.table("tickets").select("*").eq("user_id", st.session_state.user.id).order("created_at", desc=True).execute()
        saved_tickets = db_response.data
        
        if saved_tickets:
            for item in saved_tickets:
                
                status_icon = "🟢" if item.get('status') == 'Awaiting PM Scoping' else "📝"
                expander_title = f"{status_icon} Quote: {item['summary'][:55]}..."
                
                with st.expander(expander_title):
                    try:
                        past_data = json.loads(clean_json_output(item['full_data']))
                    except:
                        past_data = {}
                        
                    score = past_data.get('feasibility_score', 'Unknown')
                    hist_raw_cost = past_data.get("budget_estimate_usd", "0-0")
                    hist_low = hist_raw_cost.split("-")[0] if "-" in hist_raw_cost else hist_raw_cost
                    hist_high = hist_raw_cost.split("-")[1] if "-" in hist_raw_cost else hist_raw_cost
                    hist_fmt_low = convert_currency(hist_low, currency)
                    hist_fmt_high = convert_currency(hist_high, currency)
                    
                    ticket_name = past_data.get('project_summary', 'New Project Request')[:50] + "..."
                    sales_body = f"Hello Sales Team,\n\nPlease review the initial Sales & Feasibility scoping for the upcoming client request.\n\n-> FEASIBILITY SCORE: {score}\n-> EST. TIMELINE: {past_data.get('estimated_timeline', 'N/A')}\n-> EST. BUDGET: {hist_fmt_low} - {hist_fmt_high}\n\n-> PROJECT SUMMARY:\n{past_data.get('project_summary', 'N/A')}\n\n-> CRITICAL 'ASK' LIST (Questions for the Client):\n"
                    for q in past_data.get("client_questions", []): sales_body += f"- {q}\n"
                    sales_body += f"\n-> DEAL BREAKERS / RISKS:\n"
                    for r in past_data.get("deal_breakers", []): sales_body += f"- {r}\n"
                    sales_body += "\nBest,\nBridgeBuild Sales Hub"
                    sales_mailto = f"mailto:?subject={quote(f'Sales Quote: {ticket_name}')}&body={quote(sales_body)}"

                    current_status = item.get('status', 'Draft')
                    if current_status == 'Draft':
                        st.caption("Status: **Draft (Not Sent)**")
                    else:
                        st.caption(f"Status: **{current_status}** 🚀")
                        
                    hist_btn_col1, hist_btn_col2 = st.columns([3, 1])
                    with hist_btn_col1:
                        st.download_button("Download PDF", data=generate_local_sales_pdf(past_data, currency), file_name=f"sales_quote_{item['id'][:8]}.pdf", mime="application/pdf", key=f"hist_pdf_sales_{item['id']}", use_container_width=True)
                        
                        # --- NEW PPTX EXPORT BUTTON (HISTORY) ---
                        st.download_button("📥 Download Pitch Deck", data=generate_sales_pptx(past_data, currency), file_name=f"pitch_deck_{item['id'][:8]}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", key=f"hist_pptx_sales_{item['id']}", use_container_width=True)
                        
                    with hist_btn_col2:
                        with st.popover("Delete", use_container_width=True):
                            st.warning("Are you sure?")
                            if st.button("Confirm Delete", key=f"confirm_del_sales_{item['id']}", type="primary"):
                                try:
                                    supabase.table("tickets").delete().eq("id", item['id']).execute()
                                    st.rerun()
                                except Exception as e:
                                    st.error(f"Failed to delete: {str(e)}")
                    
                    st.markdown(f"""
                        <div style="margin-top: 5px; margin-bottom: 5px;">
                            <a href="{sales_mailto}" target="_blank" style="text-decoration: none;">
                                <button style="width: 100%; background-color: #2E7D32; color: white; border: none; padding: 8px 16px; border-radius: 6px; font-weight: bold; cursor: pointer;">
                                    Email Sales Summary
                                </button>
                            </a>
                        </div>
                        """, unsafe_allow_html=True)

                    if current_status == 'Draft':
                        if st.button("Send to PM Hub", key=f"handoff_{item['id']}", use_container_width=True):
                            supabase.table("tickets").update({"status": "Awaiting PM Scoping", "target_department": "PM"}).eq("id", item['id']).execute()
                            st.rerun()
                            
                    st.divider()
                    
                    with st.container(border=True):
                        col_m1, col_m2, col_m3 = st.columns(3)
                        with col_m1: st.metric("Feasibility", score)
                        with col_m2: st.metric("Budget", f"{hist_fmt_low} - {hist_fmt_high}") 
                        with col_m3: st.metric("Timeline", item.get('time', 'N/A'))
                    
                    st.write("")
                    col_q, col_r = st.columns(2)
                    with col_q:
                        st.markdown("##### The 'Ask' List")
                        for q in past_data.get("client_questions", []): 
                            prefix = q.split(":", 1)[0] + ":" if ":" in q else ""
                            rest = q.split(":", 1)[1] if ":" in q else q
                            st.info(f"**{prefix}** {rest.strip()}" if prefix else rest)
                    with col_r:
                        st.markdown("##### Deal Breakers")
                        for r in past_data.get("deal_breakers", []): 
                            prefix = r.split(":", 1)[0] + ":" if ":" in r else ""
                            rest = r.split(":", 1)[1] if ":" in r else r
                            st.error(f"**{prefix}** {rest.strip()}" if prefix else rest)
        else:
            st.info("No saved quotes yet. Run an analysis above!")
            
    except Exception as e:
        st.error(f"Could not load history: {str(e)}")
